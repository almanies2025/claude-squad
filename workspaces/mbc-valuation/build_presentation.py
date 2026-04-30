"""
MBC Group Valuation — PowerPoint Presentation Builder
Generates a professional 15-slide IB-style pitch deck.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
import os

OUTPUT_PATH = "/mnt/c/Users/User/Documents/GitHub/claude-squad/workspaces/mbc-valuation/MBC_Group_Presentation.pptx"

# ─── Color palette (investment banking standard) ──────────────────────────────
NAVY = RGBColor(0x1F, 0x38, 0x64)
MID_BLUE = RGBColor(0x2E, 0x75, 0xB6)
LIGHT_BLUE = RGBColor(0xBD, 0xD7, 0xEE)
GOLD = RGBColor(0xC9, 0xA2, 0x27)
DARK_GRAY = RGBColor(0x40, 0x40, 0x40)
MED_GRAY = RGBColor(0x80, 0x80, 0x80)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED = RGBColor(0xC0, 0x00, 0x00)
GREEN = RGBColor(0x00, 0x70, 0x30)
ORANGE = RGBColor(0xED, 0x7D, 0x31)


# ─── Helpers ──────────────────────────────────────────────────────────────────
def set_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(
    slide,
    text,
    left,
    top,
    width,
    height,
    font_size=12,
    bold=False,
    color=DARK_GRAY,
    align=PP_ALIGN.LEFT,
    italic=False,
    wrap=True,
):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_rect(
    slide, left, top, width, height, fill_color, line_color=None, line_width=None
):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left,
        top,
        width,
        height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_para(
    tf,
    text,
    font_size=11,
    bold=False,
    color=DARK_GRAY,
    align=PP_ALIGN.LEFT,
    italic=False,
    space_before=None,
    bullet=False,
):
    p = tf.add_paragraph()
    p.alignment = align
    if space_before:
        p.space_before = space_before
    if bullet:
        p.level = 1
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return p


def slide_footer(slide, page_num, total):
    add_textbox(
        slide,
        "MBC Group — Valuation Analysis  |  Confidential",
        Inches(0.3),
        Inches(7.1),
        Inches(6),
        Inches(0.3),
        font_size=7,
        color=MED_GRAY,
    )
    add_textbox(
        slide,
        f"{page_num} / {total}",
        Inches(9),
        Inches(7.1),
        Inches(0.7),
        Inches(0.3),
        font_size=7,
        color=MED_GRAY,
        align=PP_ALIGN.RIGHT,
    )


def section_label(slide, label, left=Inches(0.3), top=Inches(0.2)):
    add_rect(slide, left, top, Inches(2.5), Inches(0.28), MID_BLUE)
    add_textbox(
        slide,
        label,
        left + Inches(0.08),
        top + Inches(0.02),
        Inches(2.3),
        Inches(0.25),
        font_size=8,
        bold=True,
        color=WHITE,
    )


# ─── SLIDE BUILDERS ───────────────────────────────────────────────────────────


def slide_cover(prs):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_bg(slide, WHITE)

    # Dark navy header block
    add_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(4.5), NAVY)

    # Gold accent bar
    add_rect(slide, Inches(0), Inches(4.5), prs.slide_width, Inches(0.06), GOLD)

    # White bottom
    add_rect(slide, Inches(0), Inches(4.56), prs.slide_width, Inches(2.94), WHITE)

    # Company name
    add_textbox(
        slide,
        "MBC GROUP",
        Inches(0.6),
        Inches(0.6),
        Inches(9),
        Inches(0.9),
        font_size=36,
        bold=True,
        color=WHITE,
    )

    # Ticker
    add_textbox(
        slide,
        "4072.SR  |  Tadawul",
        Inches(0.6),
        Inches(1.45),
        Inches(6),
        Inches(0.4),
        font_size=14,
        color=LIGHT_BLUE,
    )

    # Title
    add_textbox(
        slide,
        "Equity Valuation Analysis",
        Inches(0.6),
        Inches(1.95),
        Inches(9),
        Inches(0.7),
        font_size=26,
        bold=False,
        color=WHITE,
    )

    add_textbox(
        slide,
        "DCF  ·  Comparable Companies  ·  Sum-of-Parts  ·  Synthesis",
        Inches(0.6),
        Inches(2.65),
        Inches(9),
        Inches(0.4),
        font_size=13,
        color=LIGHT_BLUE,
    )

    # Investment recommendation box
    add_rect(slide, Inches(0.6), Inches(3.4), Inches(5.5), Inches(0.85), MID_BLUE)
    add_textbox(
        slide,
        "INVESTMENT RATING:  BUY",
        Inches(0.7),
        Inches(3.45),
        Inches(5.3),
        Inches(0.35),
        font_size=14,
        bold=True,
        color=WHITE,
    )
    add_textbox(
        slide,
        "Target SAR 42–52  |  Central SAR 47  |  Upside 60–98%",
        Inches(0.7),
        Inches(3.8),
        Inches(5.3),
        Inches(0.3),
        font_size=11,
        color=LIGHT_BLUE,
    )

    # Price box
    add_rect(slide, Inches(6.5), Inches(3.3), Inches(3.2), Inches(1.1), GOLD)
    add_textbox(
        slide,
        "Current Price",
        Inches(6.6),
        Inches(3.35),
        Inches(3),
        Inches(0.3),
        font_size=10,
        color=NAVY,
    )
    add_textbox(
        slide,
        "SAR 26.26",
        Inches(6.6),
        Inches(3.6),
        Inches(3),
        Inches(0.5),
        font_size=22,
        bold=True,
        color=NAVY,
    )
    add_textbox(
        slide,
        "April 21, 2026",
        Inches(6.6),
        Inches(4.05),
        Inches(3),
        Inches(0.25),
        font_size=9,
        color=NAVY,
        italic=True,
    )

    # Analyst info
    add_textbox(
        slide,
        "MBA Practice Exercise  |  April 2026",
        Inches(0.6),
        Inches(5.0),
        Inches(9),
        Inches(0.3),
        font_size=10,
        color=DARK_GRAY,
    )
    add_textbox(
        slide,
        "Data: Capital IQ  |  MBC Group Annual Report (Mar 2026)",
        Inches(0.6),
        Inches(5.3),
        Inches(9),
        Inches(0.3),
        font_size=9,
        color=MED_GRAY,
        italic=True,
    )

    slide_footer(slide, 1, 15)


def slide_executive_summary(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    section_label(slide, "EXECUTIVE SUMMARY")

    add_textbox(
        slide,
        "Investment Thesis",
        Inches(0.3),
        Inches(0.65),
        Inches(9.4),
        Inches(0.4),
        font_size=16,
        bold=True,
        color=NAVY,
    )

    # Thesis box
    add_rect(slide, Inches(0.3), Inches(1.05), Inches(9.4), Inches(0.75), LIGHT_BLUE)
    add_textbox(
        slide,
        "MBC Group trades at SAR 26.26 — a 44% discount to our fair value estimate of SAR 47. "
        "Two independent valuation methods (DCF and SOTP) converge at SAR 48–50. "
        "The discount reflects conglomerate complexity and streaming execution risk, not fundamental business quality.",
        Inches(0.45),
        Inches(1.1),
        Inches(9.1),
        Inches(0.65),
        font_size=11,
        color=NAVY,
    )

    # Three segments
    add_textbox(
        slide,
        "Three Businesses, One Discount",
        Inches(0.3),
        Inches(1.95),
        Inches(9.4),
        Inches(0.35),
        font_size=13,
        bold=True,
        color=DARK_GRAY,
    )

    seg_data = [
        (
            "BOCA (Broadcast)",
            "54% of revenue",
            "MENA's dominant pan-Arab TV platform",
            "SAR 8,057M equity",
            "~SAR 6.5B at 10.5x",
        ),
        (
            "Shahid (Streaming)",
            "23% of revenue",
            "34% revenue growth — market leader",
            "SAR 7,967M equity",
            "~SAR 1.7B at 5x revenue",
        ),
        (
            "M&E (Media & Events)",
            "23% of revenue",
            "Thin margins, Vision 2030 beneficiary",
            "SAR 417M equity",
            "~SAR 0.4B at 10x",
        ),
    ]

    x_positions = [Inches(0.3), Inches(3.45), Inches(6.6)]
    colors = [NAVY, MID_BLUE, GOLD]

    for i, (seg, rev, desc, equity, val) in enumerate(seg_data):
        x = x_positions[i]
        c = colors[i]

        add_rect(slide, x, Inches(2.35), Inches(2.9), Inches(0.28), c)
        add_textbox(
            slide,
            seg,
            x + Inches(0.08),
            Inches(2.37),
            Inches(2.7),
            Inches(0.22),
            font_size=9,
            bold=True,
            color=WHITE,
        )

        add_rect(slide, x, Inches(2.65), Inches(2.9), Inches(1.8), LIGHT_GRAY)

        txb = slide.shapes.add_textbox(
            x + Inches(0.1), Inches(2.75), Inches(2.7), Inches(1.6)
        )
        tf = txb.text_frame
        tf.word_wrap = True
        add_para(tf, rev, font_size=9, color=c, bold=True)
        add_para(tf, desc, font_size=9, color=DARK_GRAY, space_before=Pt(4))
        add_para(
            tf,
            f"Equity value: {equity}",
            font_size=9,
            color=DARK_GRAY,
            space_before=Pt(6),
            bold=True,
        )
        add_para(tf, f"Implies: {val}", font_size=8, color=MED_GRAY, italic=True)

    # Valuation outputs
    add_textbox(
        slide,
        "Valuation Summary",
        Inches(0.3),
        Inches(4.6),
        Inches(9.4),
        Inches(0.35),
        font_size=13,
        bold=True,
        color=DARK_GRAY,
    )

    headers = ["Method", "Fair Value", "vs. Market", "Notes"]
    x_pos = [Inches(0.3), Inches(3.4), Inches(5.3), Inches(6.4)]
    widths = [Inches(3.0), Inches(1.8), Inches(1.0), Inches(3.3)]

    for i, (h, x, w) in enumerate(zip(headers, x_pos, widths)):
        add_rect(slide, x, Inches(4.95), w, Inches(0.3), NAVY)
        add_textbox(
            slide,
            h,
            x + Inches(0.08),
            Inches(4.97),
            w,
            Inches(0.26),
            font_size=9,
            bold=True,
            color=WHITE,
        )

    rows = [
        ("DCF Base Case (WACC 10%)", "SAR 48.8", "+86%", "Converges with SOTP"),
        ("SOTP Base Case", "SAR 49.4", "+88%", "Segment-level valuation"),
        ("Comparable Companies", "SAR 35–50", "+33–90%", "Blended peer multiples"),
        (
            "Blended Central Estimate",
            "SAR 42–52",
            "+60–98%",
            "Weighted: DCF 35%, SOTP 35%",
        ),
    ]

    for j, (m, fv, vs, n) in enumerate(rows):
        y = Inches(5.28 + j * 0.33)
        bg = LIGHT_GRAY if j % 2 == 0 else WHITE
        for i, (text, x, w) in enumerate(zip([m, fv, vs, n], x_pos, widths)):
            is_bold = j == 3
            add_rect(slide, x, y, w, Inches(0.3), bg)
            add_textbox(
                slide,
                text,
                x + Inches(0.08),
                y + Inches(0.02),
                w,
                Inches(0.26),
                font_size=9,
                bold=is_bold,
                color=(
                    GREEN
                    if ("+" in vs and j < 3)
                    else DARK_GRAY
                    if i != 2
                    else (RED if "-" in vs else DARK_GRAY)
                ),
            )

    slide_footer(slide, 2, 15)


def slide_company_overview(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    section_label(slide, "COMPANY OVERVIEW")

    # Left: Key facts
    add_textbox(
        slide,
        "MBC Group at a Glance",
        Inches(0.3),
        Inches(0.65),
        Inches(4.5),
        Inches(0.35),
        font_size=14,
        bold=True,
        color=NAVY,
    )

    facts = [
        ("Ticker:", "4072.SR — Tadawul (TASI)"),
        ("Listed:", "January 8, 2024"),
        ("Ownership:", "PIF majority; founder Waleed bin Ibrahim as chairman"),
        ("Stock price:", "SAR 26.26 (April 21, 2026)"),
        ("Market cap:", "SAR 8,844M  (~$2.36B)"),
        ("TEV:", "SAR 7,683M"),
        ("Net cash:", "SAR ~1,231M  (cash > debt)"),
        ("Shares out:", "332.5M"),
    ]

    for i, (k, v) in enumerate(facts):
        y = Inches(1.1 + i * 0.31)
        add_textbox(
            slide,
            k,
            Inches(0.35),
            y,
            Inches(1.3),
            Inches(0.28),
            font_size=9,
            bold=True,
            color=MID_BLUE,
        )
        add_textbox(
            slide,
            v,
            Inches(1.7),
            y,
            Inches(3.0),
            Inches(0.28),
            font_size=9,
            color=DARK_GRAY,
        )

    # Right: Revenue split
    add_textbox(
        slide,
        "Revenue Mix (FY2025)",
        Inches(5.3),
        Inches(0.65),
        Inches(4.5),
        Inches(0.35),
        font_size=14,
        bold=True,
        color=NAVY,
    )

    segments = [
        ("BOCA", "54%", "Pan-Arab free-to-air TV", NAVY, Inches(5.3)),
        ("Shahid", "23%", "OTT streaming", MID_BLUE, Inches(6.8)),
        ("M&E", "23%", "Production, events", GOLD, Inches(8.3)),
    ]

    for seg, pct, desc, color, x in segments:
        add_rect(slide, x, Inches(1.1), Inches(1.2), Inches(1.8), color)
        add_textbox(
            slide,
            pct,
            x,
            Inches(1.3),
            Inches(1.2),
            Inches(0.5),
            font_size=20,
            bold=True,
            color=WHITE,
            align=PP_ALIGN.CENTER,
        )
        add_textbox(
            slide,
            seg,
            x,
            Inches(1.85),
            Inches(1.2),
            Inches(0.3),
            font_size=10,
            bold=True,
            color=WHITE,
            align=PP_ALIGN.CENTER,
        )
        add_textbox(
            slide,
            desc,
            x,
            Inches(2.15),
            Inches(1.2),
            Inches(0.4),
            font_size=8,
            color=WHITE,
            align=PP_ALIGN.CENTER,
        )

    # Key metrics table
    add_textbox(
        slide,
        "Historical Financials (SAR millions)",
        Inches(0.3),
        Inches(3.7),
        Inches(9.4),
        Inches(0.35),
        font_size=13,
        bold=True,
        color=DARK_GRAY,
    )

    headers = ["", "FY2023", "FY2024", "FY2025", "FY2026E", "FY2027E"]
    col_x = [
        Inches(0.3),
        Inches(2.5),
        Inches(3.8),
        Inches(5.1),
        Inches(6.4),
        Inches(7.7),
    ]
    col_w = [
        Inches(2.1),
        Inches(1.2),
        Inches(1.2),
        Inches(1.2),
        Inches(1.2),
        Inches(1.2),
    ]

    for i, (h, x, w) in enumerate(zip(headers, col_x, col_w)):
        add_rect(slide, x, Inches(4.1), w, Inches(0.28), NAVY)
        add_textbox(
            slide,
            h,
            x + Inches(0.05),
            Inches(4.12),
            w,
            Inches(0.24),
            font_size=9,
            bold=True,
            color=WHITE,
            align=PP_ALIGN.CENTER,
        )

    rows = [
        ("Revenue", "2,559", "4,185", "5,379", "5,516", "6,088"),
        ("Revenue growth", "—", "+63.5%", "+28.5%", "+2.3%", "+10.4%"),
        ("EBITDA", "(105)", "217", "267", "693", "934"),
        ("EBITDA margin", "(4.1%)", "5.2%", "5.0%", "12.6%", "15.3%"),
        ("CapEx", "60", "198", "221", "226", "250"),
    ]

    for j, row in enumerate(rows):
        y = Inches(4.4 + j * 0.3)
        bg = LIGHT_GRAY if j % 2 == 0 else WHITE
        for i, (val, x, w) in enumerate(zip(row, col_x, col_w)):
            add_rect(slide, x, y, w, Inches(0.28), bg)
            align = PP_ALIGN.LEFT if i == 0 else PP_ALIGN.RIGHT
            add_textbox(
                slide,
                val,
                x + Inches(0.05),
                y + Inches(0.02),
                w,
                Inches(0.24),
                font_size=9,
                bold=(i == 0),
                color=DARK_GRAY,
                align=align,
            )

    # Key ratios
    add_textbox(
        slide,
        "Key Ratios",
        Inches(0.3),
        Inches(6.1),
        Inches(9.4),
        Inches(0.3),
        font_size=13,
        bold=True,
        color=DARK_GRAY,
    )

    ratios = [
        ("ROE", "8.8%", "Gross margin", "26.5%", "Debt/Equity", "3%"),
        ("EV/EBITDA", "26.9x", "EBITDA margin", "5.0%", "P/B", "1.9x"),
        (
            "P/E (TTM)",
            "22.8x",
            "Current ratio",
            "1.76x",
            "Revenue CAGR '20–'25",
            "~18%",
        ),
    ]

    for j, row in enumerate(ratios):
        y = Inches(6.45 + j * 0.28)
        for i, (k, v) in enumerate([(row[2 * k], row[2 * k + 1]) for k in range(3)]):
            x = Inches(0.3 + i * 3.15)
            add_textbox(
                slide,
                f"{row[2 * i]}:",
                x,
                y,
                Inches(1.0),
                Inches(0.25),
                font_size=8,
                bold=True,
                color=MID_BLUE,
            )
            add_textbox(
                slide,
                row[2 * i + 1],
                x + Inches(1.05),
                y,
                Inches(1.8),
                Inches(0.25),
                font_size=8,
                color=DARK_GRAY,
            )

    slide_footer(slide, 3, 15)


def slide_dcf_assumptions(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    section_label(slide, "DCF VALUATION — KEY ASSUMPTIONS")

    add_textbox(
        slide,
        "Revenue Build: FY2025–2030",
        Inches(0.3),
        Inches(0.65),
        Inches(5.5),
        Inches(0.35),
        font_size=14,
        bold=True,
        color=NAVY,
    )

    # Revenue chart
    years = ["FY25", "FY26E", "FY27E", "FY28E", "FY29E", "FY30E"]
    revenues = [5379, 5516, 6088, 6696, 7232, 7694]

    chart_data = CategoryChartData()
    chart_data.categories = years
    chart_data.add_series("Revenue (SAR M)", revenues)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.3),
        Inches(1.05),
        Inches(5.5),
        Inches(2.6),
        chart_data,
    ).chart
    chart.has_legend = False
    chart.plots[0].series[0].format.fill.solid()
    chart.plots[0].series[0].format.fill.fore_color.rgb = MID_BLUE

    # EBITDA margin progression
    add_textbox(
        slide,
        "EBITDA Margin Expansion",
        Inches(6.0),
        Inches(0.65),
        Inches(3.8),
        Inches(0.35),
        font_size=14,
        bold=True,
        color=NAVY,
    )

    margins = [
        ("FY2025", "5.0%", RED),
        ("FY2026E", "12.6%", ORANGE),
        ("FY2027E", "15.3%", ORANGE),
        ("FY2028E", "16.0%", MID_BLUE),
        ("FY2029E", "17.0%", MID_BLUE),
        ("FY2030E", "18.0%", GREEN),
    ]

    for i, (yr, margin, color) in enumerate(margins):
        y = Inches(1.1 + i * 0.42)
        add_textbox(
            slide,
            yr,
            Inches(6.1),
            y,
            Inches(1.3),
            Inches(0.3),
            font_size=9,
            color=DARK_GRAY,
        )
        bar_width = float(margin.replace("%", "")) / 18 * Inches(2.3)
        add_rect(slide, Inches(7.5), y + Inches(0.04), bar_width, Inches(0.22), color)
        add_textbox(
            slide,
            margin,
            Inches(7.5) + bar_width + Inches(0.05),
            y + Inches(0.02),
            Inches(0.5),
            Inches(0.22),
            font_size=9,
            bold=True,
            color=DARK_GRAY,
        )

    # WACC & Terminal Value
    add_textbox(
        slide,
        "Discount Rate & Terminal Value",
        Inches(0.3),
        Inches(3.8),
        Inches(9.4),
        Inches(0.35),
        font_size=14,
        bold=True,
        color=NAVY,
    )

    # WACC scenarios
    wacc_data = [
        ("Bear", "12%", "WACC 12%", RED, Inches(0.3)),
        ("Base", "10%", "WACC 10%", ORANGE, Inches(2.8)),
        ("Bull", "8%", "WACC 8%", GREEN, Inches(5.3)),
    ]

    for scenario, rate, label, color, x in wacc_data:
        add_rect(slide, x, Inches(4.2), Inches(2.3), Inches(0.75), color)
        add_textbox(
            slide,
            scenario,
            x + Inches(0.08),
            Inches(4.22),
            Inches(2.1),
            Inches(0.28),
            font_size=10,
            bold=True,
            color=WHITE,
        )
        add_textbox(
            slide,
            rate,
            x + Inches(0.08),
            Inches(4.5),
            Inches(2.1),
            Inches(0.3),
            font_size=14,
            bold=True,
            color=WHITE,
        )
        add_textbox(
            slide,
            label,
            x + Inches(0.08),
            Inches(4.78),
            Inches(2.1),
            Inches(0.2),
            font_size=8,
            color=WHITE,
            italic=True,
        )

    # Terminal value inputs
    tv_items = [
        ("Terminal growth rate:", "4.5%  (long-run Saudi nominal GDP)"),
        ("Risk-free rate:", "4.75%  (Saudi 5Y sovereign yield)"),
        ("Equity risk premium:", "5.5%  (EM media sector)"),
        ("Terminal FCFF (FY2030):", "SAR 924.8M"),
    ]

    x = Inches(0.3)
    for i, (k, v) in enumerate(tv_items):
        y = Inches(5.15 + i * 0.32)
        add_textbox(
            slide,
            k,
            x,
            y,
            Inches(2.5),
            Inches(0.28),
            font_size=9,
            bold=True,
            color=MID_BLUE,
        )
        add_textbox(
            slide,
            v,
            x + Inches(2.6),
            y,
            Inches(3.5),
            Inches(0.28),
            font_size=9,
            color=DARK_GRAY,
        )

    add_textbox(
        slide,
        "Why no beta?",
        Inches(6.5),
        Inches(5.0),
        Inches(3.3),
        Inches(0.3),
        font_size=11,
        bold=True,
        color=NAVY,
    )
    add_textbox(
        slide,
        "MBC listed January 2024. Only ~2 years of public trading history. "
        "Yahoo Finance and Capital IQ show no disclosed beta. WACC is estimated, "
        "not market-derived. Sensitivity analysis is the honest output.",
        Inches(6.5),
        Inches(5.3),
        Inches(3.3),
        Inches(1.5),
        font_size=9,
        color=DARK_GRAY,
        italic=True,
    )

    slide_footer(slide, 4, 15)


def slide_dcf_output(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    section_label(slide, "DCF VALUATION — OUTPUT")

    add_textbox(
        slide,
        "Share Price Sensitivity: WACC vs. Terminal Growth",
        Inches(0.3),
        Inches(0.65),
        Inches(9.4),
        Inches(0.35),
        font_size=14,
        bold=True,
        color=NAVY,
    )

    # Sensitivity table
    g_vals = ["3.5%", "4.0%", "4.5%", "5.0%"]
    wacc_vals = [
        ("WACC 8%", 85.3, 96.5, 112.4, 136.7),
        ("WACC 9%", 62.1, 68.3, 76.2, 86.5),
        ("WACC 10%", 48.8, 52.6, 57.3, 63.6),
        ("WACC 11%", 40.0, 42.6, 45.7, 49.4),
        ("WACC 12%", 33.8, 35.6, 37.7, 40.2),
    ]

    # Header row
    headers = ["WACC \\ g"] + g_vals
    x_start = Inches(0.5)
    col_w = Inches(1.6)

    for i, h in enumerate(headers):
        x = x_start + i * col_w
        add_rect(slide, x, Inches(1.05), col_w, Inches(0.3), NAVY)
        add_textbox(
            slide,
            h,
            x + Inches(0.05),
            Inches(1.07),
            col_w,
            Inches(0.26),
            font_size=9,
            bold=True,
            color=WHITE,
            align=PP_ALIGN.CENTER,
        )

    for row_i, (wacc_label, *prices) in enumerate(wacc_vals):
        y = Inches(1.38 + row_i * 0.36)
        is_base = wacc_label == "WACC 10%"

        add_rect(
            slide,
            x_start,
            y,
            col_w,
            Inches(0.33),
            VERY_LIGHT := RGBColor(0xDE, 0xEA, 0xF1) if is_base else LIGHT_GRAY,
        )
        add_textbox(
            slide,
            wacc_label,
            x_start + Inches(0.05),
            y + Inches(0.03),
            col_w,
            Inches(0.27),
            font_size=9,
            bold=is_base,
            color=NAVY if is_base else DARK_GRAY,
        )

        for col_i, price in enumerate(prices):
            x = x_start + (col_i + 1) * col_w
            is_bold_cell = is_base and col_i == 2
            bg = (
                RGBColor(0xCC, 0xFF, 0xCC)
                if is_bold_cell
                else (LIGHT_GRAY if row_i % 2 == 0 else WHITE)
            )
            add_rect(slide, x, y, col_w, Inches(0.33), bg)
            color = RGBColor(0x00, 0x70, 0x30) if is_bold_cell else DARK_GRAY
            add_textbox(
                slide,
                f"SAR {price:.1f}",
                x + Inches(0.05),
                y + Inches(0.03),
                col_w,
                Inches(0.27),
                font_size=9,
                bold=is_bold_cell,
                color=color,
                align=PP_ALIGN.CENTER,
            )

    # Valuation bridge
    add_textbox(
        slide,
        "Equity Bridge (Base Case: WACC 10%, g 4.5%)",
        Inches(0.3),
        Inches(3.5),
        Inches(9.4),
        Inches(0.35),
        font_size=14,
        bold=True,
        color=NAVY,
    )

    bridge_items = [
        ("PV of FCFF (FY2026–30)", "SAR 2,559M", "Sum of discounted cash flows"),
        ("PV of Terminal Value", "SAR 12,511M", "TV / (1+WACC)^5"),
        ("Enterprise Value", "SAR 15,070M", "PV FCF + PV TV"),
        ("+ Cash & ST investments", "SAR 1,352M", ""),
        ("− Total debt", "SAR (121M)", ""),
        ("− Minority interest", "SAR (70M)", ""),
        ("Equity Value", "SAR 16,231M", "Base case intrinsic value"),
        ("Shares outstanding", "332.5M", ""),
        ("Implied share price", "SAR 48.8", "+86% to market"),
    ]

    for i, (k, v, n) in enumerate(bridge_items):
        y = Inches(3.95 + i * 0.31)
        is_total = i in [2, 6]
        is_price = i == 8
        bg = (
            LIGHT_BLUE
            if is_total
            else (GREEN if is_price else (LIGHT_GRAY if i % 2 == 0 else WHITE))
        )

        add_textbox(
            slide,
            k,
            Inches(0.5),
            y,
            Inches(3.0),
            Inches(0.28),
            font_size=9,
            bold=is_total,
            color=DARK_GRAY,
        )
        add_textbox(
            slide,
            v,
            Inches(3.6),
            y,
            Inches(1.8),
            Inches(0.28),
            font_size=9,
            bold=is_total or is_price,
            color=NAVY if is_price else (GREEN if is_price else DARK_GRAY),
            align=PP_ALIGN.RIGHT,
        )
        if n:
            add_textbox(
                slide,
                n,
                Inches(5.5),
                y,
                Inches(4.0),
                Inches(0.28),
                font_size=8,
                color=MED_GRAY,
                italic=True,
            )

    slide_footer(slide, 5, 15)


def slide_comparable_companies(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    section_label(slide, "COMPARABLE COMPANIES ANALYSIS")

    add_textbox(
        slide,
        "Peer Group Multiples",
        Inches(0.3),
        Inches(0.65),
        Inches(9.4),
        Inches(0.35),
        font_size=14,
        bold=True,
        color=NAVY,
    )

    # Peer table
    peers = [
        (
            "Netflix",
            "NFLX",
            "$393B",
            "26.3x",
            "8.4x",
            "Hypergrowth streaming — global scale, FCF+",
        ),
        (
            "Disney",
            "DIS",
            "$220B",
            "20.0x",
            "2.6x",
            "Bundle strategy, parks resilience",
        ),
        (
            "Warner Bros Disc.",
            "WBD",
            "—",
            "3.1x",
            "2.6x",
            "Distressed — debt, declining linear TV",
        ),
        (
            "Paramount (deal)",
            "PARA",
            "—",
            "7.5x",
            "2.5x",
            "Skydance acquisition at 7.5x — strategic premium",
        ),
        (
            "MBC (subject)",
            "4072",
            "SAR 8.8B",
            "11.1x",
            "1.4x",
            "NTM EBITDA SAR 692.8M at SAR 26.26",
        ),
    ]

    headers = ["Company", "Ticker", "Mkt Cap", "EV/EBITDA", "EV/Rev", "Notes"]
    col_x = [
        Inches(0.3),
        Inches(2.3),
        Inches(3.2),
        Inches(4.2),
        Inches(5.2),
        Inches(6.2),
    ]
    col_w = [
        Inches(1.9),
        Inches(0.85),
        Inches(0.85),
        Inches(0.9),
        Inches(0.9),
        Inches(3.5),
    ]

    for i, (h, x, w) in enumerate(zip(headers, col_x, col_w)):
        add_rect(slide, x, Inches(1.05), w, Inches(0.28), NAVY)
        add_textbox(
            slide,
            h,
            x + Inches(0.05),
            Inches(1.07),
            w,
            Inches(0.24),
            font_size=9,
            bold=True,
            color=WHITE,
            align=PP_ALIGN.CENTER,
        )

    for j, row in enumerate(peers):
        y = Inches(1.36 + j * 0.35)
        is_mbc = j == 4
        bg = LIGHT_BLUE if is_mbc else (LIGHT_GRAY if j % 2 == 0 else WHITE)
        for i, (val, x, w) in enumerate(zip(row, col_x, col_w)):
            add_rect(slide, x, y, w, Inches(0.32), bg)
            align = PP_ALIGN.LEFT if i in [0, 5] else PP_ALIGN.CENTER
            add_textbox(
                slide,
                val,
                x + Inches(0.05),
                y + Inches(0.04),
                w,
                Inches(0.26),
                font_size=9,
                bold=is_mbc,
                color=NAVY if is_mbc else DARK_GRAY,
                align=align,
            )

    # Implied prices
    add_textbox(
        slide,
        "Implied Share Prices — Applied to MBC NTM EBITDA (SAR 692.8M)",
        Inches(0.3),
        Inches(3.25),
        Inches(9.4),
        Inches(0.35),
        font_size=14,
        bold=True,
        color=NAVY,
    )

    multiples = [
        ("Netflix 26.3x", "SAR 18,220M", "SAR 19,451M", "SAR 58.5", "+123%", True),
        ("Disney 20x", "SAR 13,856M", "SAR 15,087M", "SAR 45.4", "+73%", True),
        ("Blended 15x", "SAR 10,392M", "SAR 11,623M", "SAR 35.0", "+33%", True),
        ("Blended 12x", "SAR  8,314M", "SAR  9,545M", "SAR 28.7", "+9%", True),
        ("Paramount 7.5x", "SAR  5,196M", "SAR  5,427M", "SAR 16.3", "−38%", False),
        ("WBD 3.1x", "SAR  2,148M", "SAR    917M", "SAR  2.8", "−89%", False),
    ]

    headers2 = ["Multiple", "Implied EV", "Equity Value", "Share Price", "vs. Market"]
    col2_x = [Inches(0.3), Inches(2.8), Inches(4.5), Inches(6.0), Inches(7.5)]
    col2_w = [Inches(2.4), Inches(1.6), Inches(1.4), Inches(1.4), Inches(2.2)]

    for i, (h, x, w) in enumerate(zip(headers2, col2_x, col2_w)):
        add_rect(slide, x, Inches(3.65), w, Inches(0.28), MID_BLUE)
        add_textbox(
            slide,
            h,
            x + Inches(0.05),
            Inches(3.67),
            w,
            Inches(0.24),
            font_size=9,
            bold=True,
            color=WHITE,
            align=PP_ALIGN.CENTER,
        )

    for j, (mult, ev, eq, price, vs, positive) in enumerate(multiples):
        y = Inches(3.96 + j * 0.32)
        bg = LIGHT_GRAY if j % 2 == 0 else WHITE
        for i, (val, x, w) in enumerate(zip([mult, ev, eq, price, vs], col2_x, col2_w)):
            add_rect(slide, x, y, w, Inches(0.29), bg)
            color = (GREEN if positive else RED) if i == 4 else DARK_GRAY
            add_textbox(
                slide,
                val,
                x + Inches(0.05),
                y + Inches(0.03),
                w,
                Inches(0.24),
                font_size=9,
                bold=(i == 3),
                color=color,
                align=PP_ALIGN.LEFT if i == 0 else PP_ALIGN.CENTER,
            )

    # Key insight
    add_rect(slide, Inches(0.3), Inches(5.9), Inches(9.4), Inches(0.9), LIGHT_BLUE)
    add_textbox(
        slide,
        "Key Insight",
        Inches(0.5),
        Inches(5.95),
        Inches(2),
        Inches(0.25),
        font_size=10,
        bold=True,
        color=NAVY,
    )
    add_textbox(
        slide,
        "MBC trades at 11x NTM EBITDA while Netflix trades at 26x. The market is not yet pricing "
        "Shahid's streaming growth. At current prices, Shahid appears to be valued near zero — "
        "the market treats MBC as a mature broadcast business with a free streaming option.",
        Inches(0.5),
        Inches(6.2),
        Inches(9.0),
        Inches(0.55),
        font_size=9,
        color=NAVY,
    )

    slide_footer(slide, 6, 15)


def slide_sotp(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    section_label(slide, "SUM-OF-PARTS VALUATION")

    add_textbox(
        slide,
        "Segment-Level Equity Values (SAR millions)",
        Inches(0.3),
        Inches(0.65),
        Inches(9.4),
        Inches(0.35),
        font_size=14,
        bold=True,
        color=NAVY,
    )

    # Three segment boxes
    segs = [
        {
            "name": "BOCA",
            "subtitle": "Broadcast",
            "bear": (5580, 6245),
            "base": (7392, 8057),
            "bull": (9504, 10169),
            "multiple": "9–12x EBITDA",
            "note": "MENA's dominant pan-Arab TV. Analog: Sinclair Broadcast. Strong advertising cash flows.",
            "color": NAVY,
        },
        {
            "name": "Shahid",
            "subtitle": "Streaming",
            "bear": (4500, 4217),
            "base": (8250, 7967),
            "bull": (14400, 14117),
            "multiple": "3–8x EV/Revenue",
            "note": "34% revenue growth. MENA's leading Arabic streaming. Not yet EBITDA positive.",
            "color": MID_BLUE,
        },
        {
            "name": "M&E",
            "subtitle": "Media & Events",
            "bear": (432, 149),
            "base": (700, 417),
            "bull": (1044, 761),
            "multiple": "8–12x EBITDA",
            "note": "Production + events. Thin margins, cyclical. Benefits from Vision 2030 entertainment boom.",
            "color": GOLD,
        },
    ]

    x_positions = [Inches(0.3), Inches(3.45), Inches(6.6)]

    for i, (seg, x) in enumerate(zip(segs, x_positions)):
        # Header
        add_rect(slide, x, Inches(1.05), Inches(2.9), Inches(0.35), seg["color"])
        add_textbox(
            slide,
            seg["name"],
            x + Inches(0.08),
            Inches(1.07),
            Inches(1.5),
            Inches(0.3),
            font_size=13,
            bold=True,
            color=WHITE,
        )
        add_textbox(
            slide,
            seg["subtitle"],
            x + Inches(1.6),
            Inches(1.09),
            Inches(1.2),
            Inches(0.28),
            font_size=9,
            color=WHITE,
            italic=True,
        )

        # Values
        bg = LIGHT_GRAY
        scenarios = [
            ("Bear", seg["bear"][1], RED),
            ("Base", seg["base"][1], GREEN),
            ("Bull", seg["bull"][1], MID_BLUE),
        ]

        for j, (scenario, equity, color) in enumerate(scenarios):
            y = Inches(1.43 + j * 0.5)
            add_rect(slide, x, y, Inches(2.9), Inches(0.45), bg)
            add_textbox(
                slide,
                scenario,
                x + Inches(0.08),
                y + Inches(0.04),
                Inches(0.6),
                Inches(0.35),
                font_size=9,
                bold=True,
                color=color,
            )
            add_textbox(
                slide,
                f"SAR {equity:,}M",
                x + Inches(0.7),
                y + Inches(0.04),
                Inches(2.1),
                Inches(0.35),
                font_size=11,
                bold=True,
                color=DARK_GRAY,
            )

        # Multiple label
        add_rect(slide, x, Inches(3.0), Inches(2.9), Inches(0.28), seg["color"])
        add_textbox(
            slide,
            seg["multiple"],
            x,
            Inches(3.02),
            Inches(2.9),
            Inches(0.24),
            font_size=9,
            bold=True,
            color=WHITE,
            align=PP_ALIGN.CENTER,
        )

        # Note
        add_textbox(
            slide,
            seg["note"],
            x,
            Inches(3.35),
            Inches(2.9),
            Inches(0.7),
            font_size=8,
            color=DARK_GRAY,
            italic=True,
        )

    # SOTP Total
    add_textbox(
        slide,
        "SOTP Equity Summary",
        Inches(0.3),
        Inches(4.15),
        Inches(9.4),
        Inches(0.35),
        font_size=14,
        bold=True,
        color=NAVY,
    )

    total_data = [
        ("", "Bear", "Base", "Bull"),
        ("BOCA", "6,245", "8,057", "10,169"),
        ("Shahid", "4,217", "7,967", "14,117"),
        ("M&E", "149", "417", "761"),
        ("Total Equity (SAR M)", "10,611", "16,441", "25,047"),
        ("Share Price (SAR)", "31.9", "49.4", "75.3"),
        ("vs. Market", "+21%", "+88%", "+187%"),
    ]

    col3_x = [Inches(0.3), Inches(3.2), Inches(4.8), Inches(6.4)]
    col3_w = [Inches(2.8), Inches(1.5), Inches(1.5), Inches(1.5)]

    for row_i, row in enumerate(total_data):
        y = Inches(4.55 + row_i * 0.33)
        is_total = row_i == 4
        is_price = row_i == 5
        is_header = row_i == 0
        for i, (val, x, w) in enumerate(zip(row, col3_x, col3_w)):
            if is_header:
                bg = NAVY
                color = WHITE
            elif is_total:
                bg = GREEN
                color = WHITE
            elif is_price:
                bg = LIGHT_BLUE
                color = NAVY
            else:
                bg = LIGHT_GRAY if row_i % 2 == 0 else WHITE
                color = DARK_GRAY

            add_rect(slide, x, y, w, Inches(0.3), bg)
            add_textbox(
                slide,
                val,
                x + Inches(0.05),
                y + Inches(0.03),
                w,
                Inches(0.25),
                font_size=9,
                bold=is_total or is_header,
                color=color,
                align=PP_ALIGN.LEFT if i == 0 else PP_ALIGN.CENTER,
            )

    # Key insight
    add_rect(slide, Inches(0.3), Inches(7.0), Inches(9.4), Inches(0.4), GOLD)
    add_textbox(
        slide,
        "BOCA alone (SAR 8,057M) ≈ MBC market cap (SAR 8,844M). "
        "Shahid + M&E are essentially FREE. Market assigns near-zero value to streaming.",
        Inches(0.45),
        Inches(7.03),
        Inches(9.1),
        Inches(0.34),
        font_size=10,
        bold=True,
        color=NAVY,
    )

    slide_footer(slide, 7, 15)


def slide_synthesis(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    section_label(slide, "VALUATION SYNTHESIS")

    add_textbox(
        slide,
        "Cross-Method Comparison",
        Inches(0.3),
        Inches(0.65),
        Inches(9.4),
        Inches(0.35),
        font_size=14,
        bold=True,
        color=NAVY,
    )

    methods = [
        ("DCF — Bear", "SAR 37.7", "+43.6%", "WACC 12%, g 4.5%", RED),
        ("DCF — Base", "SAR 48.8", "+85.8%", "WACC 10%, g 4.5%", GREEN),
        ("DCF — Bull", "SAR 112.4", "+328%", "WACC 8%, g 4.5%", MID_BLUE),
        (
            "Comps — Netflix 26x",
            "SAR 58.5",
            "+123%",
            "Stretched — NFLX not comparable",
            ORANGE,
        ),
        ("Comps — Blended 15x", "SAR 35.0", "+33%", "Conservative peer set", ORANGE),
        ("Comps — Blended 12x", "SAR 28.7", "+9%", "Bearish peer set", RED),
        ("SOTP — Bear", "SAR 31.9", "+21%", "Shahid 3x rev", RED),
        ("SOTP — Base", "SAR 49.4", "+88%", "Shahid 5x rev", GREEN),
        ("SOTP — Bull", "SAR 75.3", "+187%", "Shahid 8x rev", MID_BLUE),
        ("Paramount Precedent", "SAR 16.3", "−38%", "Acquisition floor", RED),
    ]

    # Chart
    chart_data = CategoryChartData()
    chart_data.categories = [m[0] for m in methods]
    chart_data.add_series(
        "Implied Price",
        [float(m[1].replace("SAR ", "").replace(",", "")) for m in methods],
    )
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED,
        Inches(0.3),
        Inches(1.05),
        Inches(9.4),
        Inches(3.0),
        chart_data,
    ).chart

    # Color bars
    from pptx.util import Pt

    for i, (series, color) in enumerate(
        zip(chart.series[0].points, [m[4] for m in methods])
    ):
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = color

    # Market price line annotation
    add_textbox(
        slide,
        "Market Price: SAR 26.26",
        Inches(7.0),
        Inches(0.2),
        Inches(2.5),
        Inches(0.25),
        font_size=8,
        bold=True,
        color=RED,
        italic=True,
    )

    # Convergence callout
    add_rect(slide, Inches(0.3), Inches(4.2), Inches(9.4), Inches(0.8), LIGHT_BLUE)
    add_textbox(
        slide,
        "Notable: DCF (SAR 48.8) and SOTP (SAR 49.4) converge at SAR 48–50. "
        "Two completely independent methods — one projecting future cash flows, "
        "one benchmarking comparable businesses today — arriving at the same answer "
        "is a strong signal for this fair value range.",
        Inches(0.45),
        Inches(4.25),
        Inches(9.1),
        Inches(0.7),
        font_size=10,
        color=NAVY,
        bold=False,
    )

    # Weighted blend
    add_textbox(
        slide,
        "Blended Central Estimate",
        Inches(0.3),
        Inches(5.15),
        Inches(9.4),
        Inches(0.35),
        font_size=14,
        bold=True,
        color=NAVY,
    )

    blend_items = [
        ("DCF Base Case", "35%", "48.8", "17.1", MID_BLUE),
        ("SOTP Base Case", "35%", "49.4", "17.3", MID_BLUE),
        ("Comparable Companies", "20%", "35–50", "8.5", ORANGE),
        ("Paramount Precedent", "10%", "16.3", "1.6", RED),
    ]

    col4_x = [Inches(0.3), Inches(3.3), Inches(4.7), Inches(6.0), Inches(7.4)]
    col4_w = [Inches(2.9), Inches(1.3), Inches(1.2), Inches(1.3), Inches(1.3)]

    headers4 = ["Method", "Weight", "Price", "Weighted", "Color"]
    for i, (h, x, w) in enumerate(zip(headers4, col4_x, col4_w)):
        if i < 4:
            add_rect(slide, x, Inches(5.55), w, Inches(0.28), NAVY)
            add_textbox(
                slide,
                h,
                x + Inches(0.05),
                Inches(5.57),
                w,
                Inches(0.24),
                font_size=9,
                bold=True,
                color=WHITE,
                align=PP_ALIGN.CENTER,
            )

    for j, (method, wt, price, weighted, color) in enumerate(blend_items):
        y = Inches(5.86 + j * 0.32)
        bg = LIGHT_GRAY if j % 2 == 0 else WHITE
        for i, (val, x, w) in enumerate(
            zip([method, wt, price, weighted], col4_x, col4_w)
        ):
            add_rect(slide, x, y, w, Inches(0.29), bg)
            add_textbox(
                slide,
                val,
                x + Inches(0.05),
                y + Inches(0.04),
                w,
                Inches(0.24),
                font_size=9,
                bold=(i == 3),
                color=DARK_GRAY,
                align=PP_ALIGN.LEFT if i == 0 else PP_ALIGN.CENTER,
            )

    # Total
    y = Inches(5.86 + 4 * 0.32)
    add_rect(slide, Inches(0.3), y, Inches(9.4), Inches(0.35), NAVY)
    add_textbox(
        slide,
        "Blended Central Estimate",
        Inches(0.4),
        y + Inches(0.05),
        Inches(2.8),
        Inches(0.28),
        font_size=10,
        bold=True,
        color=WHITE,
    )
    add_textbox(
        slide,
        "SAR 44–47",
        Inches(6.0),
        y + Inches(0.04),
        Inches(1.2),
        Inches(0.28),
        font_size=11,
        bold=True,
        color=WHITE,
    )

    slide_footer(slide, 8, 15)


def slide_investment_recommendation(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    section_label(slide, "INVESTMENT RECOMMENDATION")

    # Big rating block
    add_rect(slide, Inches(0.3), Inches(0.65), Inches(9.4), Inches(1.2), NAVY)
    add_textbox(
        slide,
        "BUY",
        Inches(0.5),
        Inches(0.7),
        Inches(2.5),
        Inches(0.8),
        font_size=48,
        bold=True,
        color=GREEN,
    )
    add_textbox(
        slide,
        "MBC Group  |  4072.SR",
        Inches(3.0),
        Inches(0.75),
        Inches(6.5),
        Inches(0.4),
        font_size=16,
        bold=True,
        color=WHITE,
    )

    metrics = [
        ("Current Price", "SAR 26.26", "April 21, 2026"),
        ("Target Range", "SAR 42–52", "60–98% upside"),
        ("Central Target", "SAR 47", "Blended central estimate"),
        ("Market Implied", "11x NTM EBITDA", "vs. Netflix 26x, Disney 20x"),
    ]

    for i, (k, v, sub) in enumerate(metrics):
        x = Inches(0.3 + i * 2.35)
        add_rect(
            slide,
            x,
            Inches(2.0),
            Inches(2.2),
            Inches(0.75),
            MID_BLUE if i == 3 else (GREEN if i < 3 else LIGHT_GRAY),
        )
        add_textbox(
            slide,
            k,
            x + Inches(0.1),
            Inches(2.03),
            Inches(2.0),
            Inches(0.22),
            font_size=8,
            color=LIGHT_BLUE if i < 3 else DARK_GRAY,
        )
        add_textbox(
            slide,
            v,
            x + Inches(0.1),
            Inches(2.25),
            Inches(2.0),
            Inches(0.3),
            font_size=14,
            bold=True,
            color=WHITE if i < 3 else NAVY,
        )
        add_textbox(
            slide,
            sub,
            x + Inches(0.1),
            Inches(2.55),
            Inches(2.0),
            Inches(0.2),
            font_size=8,
            color=LIGHT_BLUE if i < 3 else MED_GRAY,
            italic=True,
        )

    # Investment thesis
    add_textbox(
        slide,
        "Investment Thesis",
        Inches(0.3),
        Inches(3.0),
        Inches(9.4),
        Inches(0.35),
        font_size=14,
        bold=True,
        color=NAVY,
    )

    thesis_points = [
        (
            "BOCA finances Shahid for free",
            "BOCA (54% of revenue) alone is worth SAR 8,057M — nearly equal to MBC's entire market cap of SAR 8,844M. "
            "The market is assigning zero value to Shahid and M&E, which together contributed SAR 8,384M in our SOTP base case.",
        ),
        (
            "Streaming call option at a discount",
            "Shahid is the MENA's leading Arabic streaming platform with 34% revenue growth. "
            "At current prices, investors get Shahid's growth option for free. If Shahid reaches Netflix-like margins "
            "(20%+ EBITDA), the equity value could double from current levels.",
        ),
        (
            "Two methods, one answer",
            "DCF (SAR 48.8) and SOTP (SAR 49.4) converge at SAR 48–50. This convergence — using completely "
            "independent methodologies — is a strong validation signal and substantially de-risks the estimate.",
        ),
        (
            "Strategic optionality",
            "Paramount was acquired at 7.5x EBITDA. MBC's balance sheet (net cash SAR 1.2B) is far stronger than "
            "Paramount's. A strategic buyer (PIF or regional media group) could justify SAR 35–50 in an acquisition.",
        ),
    ]

    for i, (head, body) in enumerate(thesis_points):
        y = Inches(3.4 + i * 0.78)
        add_rect(
            slide,
            Inches(0.3),
            y,
            Inches(0.25),
            Inches(0.65),
            MID_BLUE if i % 2 == 0 else NAVY,
        )
        add_textbox(
            slide,
            head,
            Inches(0.7),
            y,
            Inches(8.9),
            Inches(0.25),
            font_size=10,
            bold=True,
            color=NAVY,
        )
        add_textbox(
            slide,
            body,
            Inches(0.7),
            y + Inches(0.25),
            Inches(8.9),
            Inches(0.45),
            font_size=8.5,
            color=DARK_GRAY,
        )

    slide_footer(slide, 9, 15)


def slide_risks(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    section_label(slide, "KEY RISKS")

    add_textbox(
        slide,
        "Risk Matrix",
        Inches(0.3),
        Inches(0.65),
        Inches(9.4),
        Inches(0.35),
        font_size=14,
        bold=True,
        color=NAVY,
    )

    risks = [
        (
            "HIGH",
            "Streaming margin expansion slower than modeled",
            "Capital IQ consensus projects FY2026 EBITDA margin of 12.6% vs. FY2025 actual 5.0%. "
            "If Shahid content costs remain elevated, margin expansion stalls and DCF/SOTP estimates compress significantly.",
            "Medium",
            "Base case DCF, SOTP both assume Shahid scale economics; compression = −30–40% to valuation",
        ),
        (
            "HIGH",
            "Linear TV advertising decline accelerates",
            "BOCA (54% of revenue) depends on pan-Arab TV advertising. Global secular decline in linear TV viewing "
            "may accelerate in MENA faster than expected, compressing BOCA's cash generation.",
            "Low–Medium",
            "BOCA generates SAR 650M+ EBITDA; even at 15% decline, still cash generative",
        ),
        (
            "MEDIUM",
            "Netflix / Disney+ MENA expansion intensifies",
            "Netflix has ~4M MENA subscribers; Disney+ is expanding. Aggressive pricing or content investment "
            "by global players could slow Shahid's subscriber growth and margin path.",
            "Medium",
            "Shahid's Arabic-language moat (local content, Arabic UI) provides some protection",
        ),
        (
            "HIGH",
            "FY2026 EBITDA misses Capital IQ consensus",
            "Consensus estimates SAR 692.8M EBITDA (12.6% margin). If actual comes in at SAR 400–500M "
            "(8–9% margin), the market re-rates lower and valuation multiples compress.",
            "Medium",
            "Capital IQ consensus from 8 brokers; track record unknown for newly listed company",
        ),
        (
            "MEDIUM",
            "Working capital volatility obscures FCF",
            "FY2025 had SAR −180.7M levered FCF despite positive EBITDA — driven by SAR 550M AR build. "
            "This pattern may repeat, creating FCF volatility that doesn't reflect business quality.",
            "Low",
            "Net cash position (SAR 1.2B) provides cushion; one-time content investment cycle",
        ),
        (
            "MEDIUM",
            "Saudi FX / currency risk",
            "All MBC revenues are in SAR. USD/SAR peg is stable, but any devaluation risk "
            "would affect USD-denominated investor returns.",
            "Very Low",
            "Peg has been stable for 40+ years; not a meaningful near-term risk",
        ),
    ]

    for i, (severity, title, body, likelihood, mitigation) in enumerate(risks):
        y = Inches(1.05 + i * 1.05)
        sev_color = RED if severity == "HIGH" else ORANGE
        add_rect(slide, Inches(0.3), y, Inches(0.55), Inches(0.9), sev_color)
        add_textbox(
            slide,
            severity,
            Inches(0.32),
            y + Inches(0.3),
            Inches(0.5),
            Inches(0.28),
            font_size=7,
            bold=True,
            color=WHITE,
            align=PP_ALIGN.CENTER,
        )

        add_textbox(
            slide,
            title,
            Inches(1.0),
            y,
            Inches(8.7),
            Inches(0.25),
            font_size=10,
            bold=True,
            color=DARK_GRAY,
        )
        add_textbox(
            slide,
            body,
            Inches(1.0),
            y + Inches(0.25),
            Inches(5.8),
            Inches(0.5),
            font_size=8,
            color=DARK_GRAY,
        )
        add_textbox(
            slide,
            f"Likelihood: {likelihood}",
            Inches(6.9),
            y + Inches(0.25),
            Inches(2.8),
            Inches(0.22),
            font_size=8,
            bold=True,
            color=MED_GRAY,
        )
        add_textbox(
            slide,
            mitigation,
            Inches(6.9),
            y + Inches(0.48),
            Inches(2.8),
            Inches(0.4),
            font_size=7.5,
            color=GREEN,
            italic=True,
        )

    slide_footer(slide, 10, 15)


def slide_valuation_floor(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    section_label(slide, "VALUATION FLOOR & BULL CASE")

    add_textbox(
        slide,
        "Paramount Precedent — Acquisition Floor",
        Inches(0.3),
        Inches(0.65),
        Inches(9.4),
        Inches(0.35),
        font_size=14,
        bold=True,
        color=NAVY,
    )

    # Paramount deal context
    add_rect(slide, Inches(0.3), Inches(1.05), Inches(9.4), Inches(1.1), LIGHT_GRAY)
    add_textbox(
        slide,
        "Skydance / National Media acquired Paramount Global at 7.5x fully synergized 2026E EBITDA "
        "(February 2026). Deal terms: control premium, synergy assumptions. "
        "MBC's balance sheet (net cash SAR 1.2B) is far stronger than Paramount's (net debt). "
        "MBC should deserve a higher multiple than Paramount as a standalone public market asset.",
        Inches(0.45),
        Inches(1.1),
        Inches(9.1),
        Inches(1.0),
        font_size=10,
        color=DARK_GRAY,
    )

    # Floor table
    floor_data = [
        ("7.5x (Paramount)", "SAR 5,196M", "SAR 5,427M", "SAR 16.3", "−38%"),
        ("10x", "SAR 6,928M", "SAR 7,159M", "SAR 21.5", "−18%"),
        ("12x", "SAR 8,314M", "SAR 9,545M", "SAR 28.7", "+9%"),
    ]

    headers5 = ["Multiple", "Applied EV", "Equity Value", "Share Price", "vs. Market"]
    col5_x = [Inches(0.3), Inches(2.8), Inches(4.7), Inches(6.2), Inches(7.7)]
    col5_w = [Inches(2.4), Inches(1.8), Inches(1.4), Inches(1.4), Inches(2.0)]

    for i, (h, x, w) in enumerate(zip(headers5, col5_x, col5_w)):
        add_rect(slide, x, Inches(2.25), w, Inches(0.28), MID_BLUE)
        add_textbox(
            slide,
            h,
            x + Inches(0.05),
            Inches(2.27),
            w,
            Inches(0.24),
            font_size=9,
            bold=True,
            color=WHITE,
            align=PP_ALIGN.CENTER,
        )

    for j, row in enumerate(floor_data):
        y = Inches(2.56 + j * 0.32)
        bg = LIGHT_GRAY if j % 2 == 0 else WHITE
        for i, (val, x, w) in enumerate(zip(row, col5_x, col5_w)):
            add_rect(slide, x, y, w, Inches(0.29), bg)
            color = RED if "-" in row[4] else DARK_GRAY
            add_textbox(
                slide,
                val,
                x + Inches(0.05),
                y + Inches(0.03),
                w,
                Inches(0.24),
                font_size=9,
                bold=(i == 3),
                color=color,
                align=PP_ALIGN.LEFT if i == 0 else PP_ALIGN.CENTER,
            )

    # Bull case drivers
    add_textbox(
        slide,
        "Bull Case — Path to SAR 65–75",
        Inches(0.3),
        Inches(3.4),
        Inches(9.4),
        Inches(0.35),
        font_size=14,
        bold=True,
        color=NAVY,
    )

    bull_points = [
        (
            "Shahid reaches Netflix-like scale",
            "5M+ subscribers, 25%+ EBITDA margins, valued at 8x+ EV/Revenue → Shahid equity SAR 14,000M+",
        ),
        (
            "BOCA holds margins at 22%+",
            "Saudi advertising market stays robust; BOCA remains dominant platform; SAR 9,500M+ equity",
        ),
        (
            "DCF terminal growth 5%",
            "Saudi digitalization premium embedded; WACC 8% (country risk premium compresses)",
        ),
        (
            "International analyst coverage initiates",
            "Coverage by MSCI EM media analysts raises visibility; institutional ownership increases → multiple re-rating",
        ),
        (
            "Strategic transaction",
            "PIF or regional media group acquires MBC at control premium (10–12x EBITDA) → SAR 50–65 per share",
        ),
    ]

    for i, (head, body) in enumerate(bull_points):
        y = Inches(3.85 + i * 0.58)
        add_rect(slide, Inches(0.3), y, Inches(0.25), Inches(0.5), GREEN)
        add_textbox(
            slide,
            head,
            Inches(0.7),
            y,
            Inches(8.9),
            Inches(0.25),
            font_size=10,
            bold=True,
            color=DARK_GRAY,
        )
        add_textbox(
            slide,
            body,
            Inches(0.7),
            y + Inches(0.25),
            Inches(8.9),
            Inches(0.3),
            font_size=9,
            color=MED_GRAY,
        )

    slide_footer(slide, 11, 15)


def slide_discussion(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    section_label(slide, "DISCUSSION QUESTIONS")

    add_textbox(
        slide,
        "Three Questions for the Class",
        Inches(0.3),
        Inches(0.65),
        Inches(9.4),
        Inches(0.35),
        font_size=14,
        bold=True,
        color=NAVY,
    )

    questions = [
        (
            "Q1",
            "The market prices MBC at 11x NTM EBITDA while Netflix trades at 26x. "
            "Does that gap reflect genuine business quality differences — or does it reflect that MBC isn't "
            "widely covered by international analysts yet (it listed January 2024)?",
            "What would it take for MBC to re-rate to Netflix-like multiples? Is analyst coverage the missing "
            "catalyst, or is there a structural reason the market will never price MBC at 26x?",
        ),
        (
            "Q2",
            "Our DCF (SAR 48.8) and SOTP (SAR 49.4) converge almost exactly. "
            "Is this convergence meaningful — does it give you more confidence in the SAR 48–50 fair value range? "
            "Or could both methods be systematically wrong in the same direction?",
            "Both methods assume Shahid reaches streaming scale. What single assumption, if wrong, "
            "would most compress the valuation? Could that assumption be validated or invalidated "
            "with one more year of data?",
        ),
        (
            "Q3",
            "The Paramount acquisition happened at 7.5x EBITDA. MBC is in a stronger financial position "
            "(net cash SAR 1.2B vs. Paramount's net debt). Should MBC deserve a higher multiple than 7.5x "
            "as a public market standalone — or does the strategic premium in that deal mean the public market "
            "simply doesn't care about fundamental quality?",
            "If you were advising MBC's board: would you recommend a share buyback program, "
            "a strategic acquisition, or maintaining the net cash balance as a war chest?",
        ),
    ]

    for i, (qnum, question, probe) in enumerate(questions):
        y = Inches(1.1 + i * 2.0)
        add_rect(slide, Inches(0.3), y, Inches(0.55), Inches(1.75), NAVY)
        add_textbox(
            slide,
            qnum,
            Inches(0.32),
            y + Inches(0.6),
            Inches(0.5),
            Inches(0.5),
            font_size=16,
            bold=True,
            color=WHITE,
            align=PP_ALIGN.CENTER,
        )

        add_textbox(
            slide,
            question,
            Inches(1.0),
            y,
            Inches(8.7),
            Inches(1.0),
            font_size=11,
            bold=True,
            color=DARK_GRAY,
        )
        add_textbox(
            slide,
            probe,
            Inches(1.0),
            y + Inches(1.05),
            Inches(8.7),
            Inches(0.65),
            font_size=9,
            color=MED_GRAY,
            italic=True,
        )

    slide_footer(slide, 12, 15)


def slide_appendix(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    section_label(slide, "APPENDIX")

    add_textbox(
        slide,
        "Data Sources & Methodology Notes",
        Inches(0.3),
        Inches(0.65),
        Inches(9.4),
        Inches(0.35),
        font_size=14,
        bold=True,
        color=NAVY,
    )

    sources = [
        (
            "Primary Data",
            "Capital IQ Excel (MBC Group SASE 4072 Financials.xls) — 14 sheets of standardized data including "
            "income statement (FY2020–2025), balance sheet, cash flow, key stats, multiples, ratios, capital structure, segments. "
            "Capital IQ consensus estimates for FY2026E and FY2027E.",
        ),
        (
            "Annual Report",
            "MBC Group Form Annual Report (March 29, 2026) — segment revenue breakdown, operational review, "
            "D&A schedule, content investment commentary. Multi-column PDF layout defeated text extraction for "
            "primary financial statements; Capital IQ Excel provided complete structured data.",
        ),
        (
            "Market Data",
            "Yahoo Finance, Saudi Exchange (Tadawul) — stock price April 21, 2026 (SAR 26.26), shares outstanding "
            "(332.5M), market cap (SAR 8,844M), TEV (SAR 7,683M). No beta available on Yahoo Finance or Capital IQ.",
        ),
        (
            "Peer Multiples",
            "Yahoo Finance, multiples.vc, Trefis, valueinvesting.io — Netflix (NFLX), Disney (DIS), "
            "Warner Bros Discovery (WBD), Paramount (PARA) trading multiples as of April 2026. "
            "Paramount deal price (Skydance/National Media, February 2026) from Paramount press release.",
        ),
        (
            "Key Methodology Choices",
            "WACC estimated (no reliable beta): 10% base case, 8% bull, 12% bear. "
            "Terminal growth: 4.5% (Saudi nominal GDP proxy). "
            "CapEx: 4.1% of revenue (FY2025 actual). "
            "Tax rate: 20% (Saudi corporate tax). "
            "FCF normalization: NWC changes excluded (working capital swings transitory). "
            "Segment EBITDA estimated: no segment-level EBITDA disclosure in available data.",
        ),
    ]

    for i, (head, body) in enumerate(sources):
        y = Inches(1.05 + i * 1.2)
        add_rect(slide, Inches(0.3), y, Inches(1.8), Inches(1.05), MID_BLUE)
        add_textbox(
            slide,
            head,
            Inches(0.35),
            y + Inches(0.3),
            Inches(1.7),
            Inches(0.4),
            font_size=9,
            bold=True,
            color=WHITE,
        )
        add_textbox(
            slide,
            body,
            Inches(2.2),
            y,
            Inches(7.5),
            Inches(1.05),
            font_size=8.5,
            color=DARK_GRAY,
        )

    slide_footer(slide, 13, 15)


def slide_disclaimer(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    section_label(slide, "IMPORTANT DISCLOSURES")

    add_rect(slide, Inches(0.3), Inches(0.65), Inches(9.4), Inches(2.0), LIGHT_GRAY)
    add_textbox(
        slide,
        "This presentation was produced as an MBA practice exercise. It is not investment advice, "
        "a solicitation, or a recommendation to buy or sell any security. All data presented herein "
        "is sourced from publicly available information including Capital IQ, the company's annual report, "
        "and market data as of April 21, 2026.",
        Inches(0.45),
        Inches(0.7),
        Inches(9.1),
        Inches(0.85),
        font_size=10,
        color=DARK_GRAY,
    )
    add_textbox(
        slide,
        "Valuations involve estimates and assumptions. DCF and SOTP models are highly sensitive to input "
        "assumptions including WACC, terminal growth rate, peer multiples, and segment margin estimates. "
        "MBC listed January 2024 — limited public trading history means no reliable beta is available. "
        "WACC is estimated, not market-derived. Actual results may differ materially from projections.",
        Inches(0.45),
        Inches(1.55),
        Inches(9.1),
        Inches(0.85),
        font_size=10,
        color=DARK_GRAY,
    )

    add_textbox(
        slide,
        "Methodology Summary",
        Inches(0.3),
        Inches(2.85),
        Inches(9.4),
        Inches(0.35),
        font_size=14,
        bold=True,
        color=NAVY,
    )

    methods_summary = [
        (
            "DCF (Discounted Cash Flow)",
            "5-year explicit FCFF projections (FY2026–2030). Revenue from Capital IQ consensus (FY2026–27) then faded "
            "to Saudi nominal GDP growth. EBITDA margins 5.0% → 18.0% (FY2025 → FY2030). Terminal value via Gordon "
            "Growth model at 4.5% perpetual growth. Three WACC scenarios (8%, 10%, 12%).",
        ),
        (
            "Comparable Companies",
            "Three peer groups: Global Streaming (Netflix, Disney), Legacy/Distressed Media (WBD, Paramount), MBC itself. "
            "Multiples applied to NTM EBITDA (SAR 692.8M) and NTM Revenue (SAR 5,515.8M). Blended peer multiple "
            "construction (12–18x NTM EBITDA) based on qualitative assessment of business quality differentials.",
        ),
        (
            "Sum-of-Parts (SOTP)",
            "Segment-level valuation: BOCA at 9–12x EBITDA (analog: regional TV broadcasters), Shahid at 3–8x EV/Revenue "
            "(analog: streaming peers), M&E at 8–12x EBITDA (analog: production/events companies). Net debt allocated "
            "pro-rata by segment revenue weight.",
        ),
        (
            "Precedent Transactions",
            "Paramount/Skydance acquisition (February 2026) at 7.5x fully synergized 2026E EBITDA used as acquisition "
            "floor. Applied to MBC NTM EBITDA (SAR 692.8M). Strategic premium acknowledged but not quantified.",
        ),
    ]

    for i, (head, body) in enumerate(methods_summary):
        y = Inches(3.25 + i * 0.9)
        add_rect(slide, Inches(0.3), y, Inches(0.25), Inches(0.75), MID_BLUE)
        add_textbox(
            slide,
            head,
            Inches(0.7),
            y,
            Inches(8.9),
            Inches(0.28),
            font_size=10,
            bold=True,
            color=NAVY,
        )
        add_textbox(
            slide,
            body,
            Inches(0.7),
            y + Inches(0.28),
            Inches(8.9),
            Inches(0.6),
            font_size=8.5,
            color=DARK_GRAY,
        )

    slide_footer(slide, 14, 15)


# ─── MAIN ──────────────────────────────────────────────────────────────────────
def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slide_cover(prs)
    slide_executive_summary(prs)
    slide_company_overview(prs)
    slide_dcf_assumptions(prs)
    slide_dcf_output(prs)
    slide_comparable_companies(prs)
    slide_sotp(prs)
    slide_synthesis(prs)
    slide_investment_recommendation(prs)
    slide_risks(prs)
    slide_valuation_floor(prs)
    slide_discussion(prs)
    slide_appendix(prs)
    slide_disclaimer(prs)

    prs.save(OUTPUT_PATH)
    print(f"Saved: {OUTPUT_PATH}")


if __name__ == "__main__":
    main()
